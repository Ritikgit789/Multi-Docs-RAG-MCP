import os
import fitz  # PyMuPDF
import pandas as pd
import docx
from pptx import Presentation
from core_mcp.mcp import MCPMessage

# Optional: use tiktoken for better chunking
try:
    import tiktoken
    HAS_TIKTOKEN = True
except ImportError:
    HAS_TIKTOKEN = False

def read_pdf(path):
    doc = fitz.open(path)
    return "\n".join([page.get_text() for page in doc])

def read_docx(path):
    doc = docx.Document(path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def read_csv(path):
    df = pd.read_csv(path)
    return df

def chunk_csv_dataframe(df, rows_per_chunk=20):
    chunks = []
    for start in range(0, len(df), rows_per_chunk):
        chunk_df = df.iloc[start:start+rows_per_chunk]
        chunk_str = chunk_df.to_string(index=False)
        if chunk_str.strip():
            chunks.append(chunk_str)
    return chunks

def read_pptx(path):
    prs = Presentation(path)
    slides = []
    for slide_num, slide in enumerate(prs.slides, 1):
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text.strip())
        if content:
            slides.append(f"Slide {slide_num}: " + " ".join(content))
    return "\n".join(slides)

def read_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()

def split_into_chunks_tiktoken(text, max_tokens=300, overlap=50):
    enc = tiktoken.get_encoding("cl100k_base")
    tokens = enc.encode(text)
    chunks = []
    start = 0
    while start < len(tokens):
        end = min(start + max_tokens, len(tokens))
        chunk = enc.decode(tokens[start:end])
        if chunk.strip():
            chunks.append(chunk.strip())
        start += max_tokens - overlap
    return chunks

def split_into_chunks_line(text, max_tokens=300):
    lines = text.split("\n")
    chunks, current = [], ""
    for line in lines:
        if len(current) + len(line) < max_tokens:
            current += line.strip() + " "
        else:
            if current.strip():
                chunks.append(current.strip())
            current = line.strip() + " "
    if current.strip():
        chunks.append(current.strip())
    return chunks

def get_text_from_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".pdf":
        return read_pdf(file_path)
    elif ext == ".docx":
        return read_docx(file_path)
    elif ext == ".csv":
        # Return DataFrame for CSV
        return read_csv(file_path)
    elif ext == ".pptx":
        return read_pptx(file_path)
    elif ext in [".txt", ".md"]:
        return read_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def run_ingestion_agent(file_path, receiver_agent="IndexAgent"):
    try:
        ext = os.path.splitext(file_path)[-1].lower()
        if ext == ".csv":
            df = get_text_from_file(file_path)
            chunks = chunk_csv_dataframe(df, rows_per_chunk=20)
        else:
            raw_text = get_text_from_file(file_path)
            if HAS_TIKTOKEN:
                chunks = split_into_chunks_tiktoken(raw_text)
            else:
                chunks = split_into_chunks_line(raw_text)

        # Placeholder for embedding generation — replace with Gemini/Groq later
        embeddings = [None] * len(chunks)

        return MCPMessage(
            sender="IngestionAgent",
            receiver=receiver_agent,
            msg_type=MCPMessage.TYPE_INDEX,
            payload={
                "chunks": chunks,
                "embeddings": embeddings,
                "source_file": os.path.basename(file_path)
            }
        )

    except Exception as e:
        return MCPMessage(
            sender="IngestionAgent",
            receiver=receiver_agent,
            msg_type=MCPMessage.TYPE_ERROR,
            payload={"error": str(e)}
        )
